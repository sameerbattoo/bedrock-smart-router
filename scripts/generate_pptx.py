"""Generate PowerPoint presentation from the HTML presentation content."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Dark theme colors
BG_COLOR = RGBColor(0x0a, 0x0e, 0x17)
ORANGE = RGBColor(0xf9, 0x73, 0x16)
LIGHT_ORANGE = RGBColor(0xfb, 0x92, 0x3c)
WHITE = RGBColor(0xf3, 0xf4, 0xf6)
GRAY = RGBColor(0x9c, 0xa3, 0xaf)
DARK_GRAY = RGBColor(0x4b, 0x55, 0x63)
CARD_BG = RGBColor(0x11, 0x18, 0x27)
GREEN = RGBColor(0x4a, 0xde, 0x80)
RED = RGBColor(0xf8, 0x71, 0x71)
BLUE = RGBColor(0x60, 0xa5, 0xfa)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

def add_card(slide, left, top, width, height, icon, title, desc):
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = RGBColor(0x1f, 0x29, 0x37)
    shape.line.width = Pt(1)
    
    # Icon + Title + Description
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = f"{icon}  {title}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(6)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Problem Statement
# ═══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide1)

add_text_box(slide1, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
             "The Model Selection Problem", font_size=36, bold=True, color=ORANGE)
add_text_box(slide1, Inches(0.8), Inches(1.3), Inches(10), Inches(0.8),
             "Amazon Bedrock offers 65+ models across 8 families. Picking the right one for each request\nis a manual, static decision that leaves money and performance on the table.",
             font_size=14, color=GRAY)

cards = [
    ("💸", "Overspending", "Using Opus for simple tasks costs 50× more than needed"),
    ("🐌", "Unnecessary Latency", "Large models add seconds for tasks a smaller model handles in ms"),
    ("💀", "No Resilience", "Hardcoded model IDs mean throttling = downtime"),
    ("🔒", "Single Family Lock-in", "Bedrock native router only routes within one family"),
    ("👥", "No Multi-Tenant Control", "Free and enterprise users hit the same model"),
    ("📊", "Blind Decisions", "No visibility into why a model was chosen"),
]

for i, (icon, title, desc) in enumerate(cards):
    col = i % 3
    row = i // 3
    left = Inches(0.8 + col * 4.1)
    top = Inches(2.5 + row * 2.3)
    add_card(slide1, left, top, Inches(3.8), Inches(2.0), icon, title, desc)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: Available Solutions
# ═══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2)

add_text_box(slide2, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Existing Solutions Fall Short", font_size=36, bold=True, color=ORANGE)
add_text_box(slide2, Inches(0.8), Inches(1.1), Inches(11), Inches(0.6),
             "Generic LLM gateways and Bedrock's native router each solve part of the problem — none address Bedrock-specific needs holistically.",
             font_size=13, color=GRAY)

# Comparison table
from pptx.util import Cm
table_data = [
    ["Capability", "LiteLLM", "OpenRouter", "Portkey", "Bedrock Native"],
    ["Purpose-built for Bedrock", "No", "No", "No", "Yes"],
    ["Cross-family routing", "Generic", "Generic", "Generic", "Single family"],
    ["CRIS profile awareness", "No", "No", "No", "Manual"],
    ["Prompt cache optimization", "No", "No", "No", "No"],
    ["Circuit breaker + fallback", "No", "No", "Yes", "No"],
    ["Multi-tenant cost tracking", "No", "No", "No", "Manual"],
    ["Quality-aware routing", "No", "No", "No", "No"],
    ["Strands Agents integration", "No", "No", "No", "No"],
    ["Zero external dependencies", "Redis", "SaaS", "SaaS", "Yes"],
]

rows, cols = len(table_data), len(table_data[0])
tbl_shape = slide2.shapes.add_table(rows, cols, Inches(0.6), Inches(1.8), Inches(12), Inches(4.8))
tbl = tbl_shape.table

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = table_data[r][c]
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_COLOR if r > 0 else RGBColor(0x11, 0x18, 0x27)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(9)
        if r == 0:
            p.font.bold = True
            p.font.color.rgb = GRAY
        elif table_data[r][c] == "Yes":
            p.font.color.rgb = GREEN
            p.font.bold = True
        elif table_data[r][c] == "No":
            p.font.color.rgb = DARK_GRAY
        else:
            p.font.color.rgb = RGBColor(0xfb, 0xbf, 0x24)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: What is Smart Router
# ═══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3)

add_text_box(slide3, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Bedrock Smart Router", font_size=36, bold=True, color=ORANGE)
add_text_box(slide3, Inches(0.8), Inches(1.1), Inches(11), Inches(0.6),
             "Drop-in replacement for converse() & converse_stream() with support for Strands Agents SDK.",
             font_size=14, color=GRAY)

features = [
    ("🧠", "Intelligent Complexity Detection", "15-dimension classifier analyzes each request in <1ms"),
    ("⚡", "4 Routing Strategies", "Cost, latency, quality, or balanced — one-word presets"),
    ("🛡️", "Built-in Resilience", "Circuit breakers, retry with backoff, multi-model fallback"),
    ("☁️", "Bedrock-Native", "CRIS, latency optimization, prompt caching, guardrails, AIPs"),
]

for i, (icon, title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(2.0 + row * 2.2)
    add_card(slide3, left, top, Inches(5.8), Inches(1.9), icon, title, desc)

# Flow steps
flow_labels = ["Analyze", "Score", "Invoke", "Fallback", "Record"]
for i, label in enumerate(flow_labels):
    left = Inches(1.0 + i * 2.4)
    shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(6.3), Inches(1.8), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = RGBColor(0x1f, 0x29, 0x37)
    tf = shape.text_frame
    tf.paragraphs[0].text = f"{i+1}. {label}"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = ORANGE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: Problem Solved
# ═══════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4)

add_text_box(slide4, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Problem → Solved", font_size=36, bold=True, color=ORANGE)
add_text_box(slide4, Inches(0.8), Inches(1.1), Inches(11), Inches(0.6),
             "Analyzes each request's complexity in real-time and routes to the optimal model — balancing cost, latency, and quality automatically.",
             font_size=13, color=GRAY)

solutions = [
    ("💰", "40-90% Cost Reduction", "Simple requests auto-route to Nova Micro instead of Opus"),
    ("⚡", "3-10× Faster", "Micro models respond in 200ms vs 2000ms for simple tasks"),
    ("🛡️", "Zero Downtime", "Circuit breakers + automatic fallback. Users never see a 429"),
    ("🌐", "Cross-Family", "Routes across Anthropic, Nova, Meta, Mistral, OpenAI"),
    ("🏢", "Per-Tenant Control", "Enterprise gets Opus, Free gets Micro. Per-tenant budgets + AIPs"),
    ("🔍", "Full Observability", "Explain every decision. A/B test. CloudWatch + OTel"),
]

for i, (icon, title, desc) in enumerate(solutions):
    col = i % 3
    row = i // 3
    left = Inches(0.8 + col * 4.1)
    top = Inches(2.0 + row * 2.3)
    add_card(slide4, left, top, Inches(3.8), Inches(2.0), icon, title, desc)

# Footer badges
add_text_box(slide4, Inches(1.5), Inches(6.3), Inches(10), Inches(0.4),
             "Drop-in replacement (1 line change)  •  Zero external dependencies (boto3 only)  •  Works with Strands Agents SDK",
             font_size=11, color=LIGHT_ORANGE, align=PP_ALIGN.CENTER)

# Contact & repo
add_text_box(slide4, Inches(1.5), Inches(6.7), Inches(10), Inches(0.3),
             "📧 Sameer Battoo (sbattoo@amazon.com)  •  🔗 https://gitlab.aws.dev/sbattoo/bedrock_smart_router",
             font_size=10, color=GRAY, align=PP_ALIGN.CENTER)

# Save
output_path = "/Users/sbattoo/Src/bedrock-smart-router/docs/Bedrock_Smart_Router.pptx"
prs.save(output_path)
print(f"✅ PowerPoint saved to: {output_path}")
